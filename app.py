from __future__ import annotations

import re
import unicodedata
import tempfile
from pathlib import Path
from typing import Dict, List, Tuple

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import streamlit as st

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None

APP_TITLE = "Dashboard Pozos Swab Lote X"

# Streamlit Cloud corre en Linux y distingue mayúsculas y minúsculas.
# Por eso se aceptan data, Data o DATA.
DATA_DIR_CANDIDATES = [
    Path(__file__).parent / "data",
    Path(__file__).parent / "Data",
    Path(__file__).parent / "DATA",
]
DATA_DIR = next((p for p in DATA_DIR_CANDIDATES if p.exists()), DATA_DIR_CANDIDATES[0])

MESES = {
    "ENERO": 1,
    "FEBRERO": 2,
    "MARZO": 3,
    "ABRIL": 4,
    "MAYO": 5,
    "JUNIO": 6,
    "JULIO": 7,
    "AGOSTO": 8,
    "SETIEMBRE": 9,
    "SEPTIEMBRE": 9,
    "OCTUBRE": 10,
    "NOVIEMBRE": 11,
    "DICIEMBRE": 12,
}

MESES_INV = {
    1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril", 5: "Mayo", 6: "Junio",
    7: "Julio", 8: "Agosto", 9: "Setiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
}

# ============================================================
# REGLA DE PERIODO
# ============================================================
# Este dashboard toma el mes y el año ÚNICAMENTE desde el nombre del archivo.
# Esto evita errores cuando una hoja interna trae títulos copiados de otro mes.
# Ejemplos válidos:
#   Estado de Pozos - Abril 2024.xlsx
#   Estado de Pozos - Abril 2025 OIG.xlsx
#   Estado de Pozos - Abril 2026 OIG.xlsx
#   Estado de Pozos - Julio OIG 2024.xlsx
# Si el nombre no tiene mes y año, el archivo se omite y se muestra en el log.

REQ = ["*PFORMACION", "*ESTADO", "*TIPO_DE_POZO", "*ULT_EST", "*BATERIA"]

# ============================================================
# POZOS CANDIDATOS ATA EXCLUIDOS DEL ANÁLISIS SWAB
# ============================================================
# Estos 88 pozos no se consideran en KPIs, tablas, gráficas ni PPT.
POZOS_ATA_EXCLUIR = [
    "AA37", "AA54", "AA76", "AA112", "AA1577", "AA1598", "AA1599", "AA1633",
    "AA1661", "AA1847", "AA1930", "AA5631", "AA5707", "AA5861", "AA5926",
    "AA5971", "AA6192", "AA6338", "AA6342", "AA6372", "AA6423", "AA6454",
    "AA6517", "AA6646", "AA6762", "AA7201", "AA9154", "AA9329", "AA9364",
    "AA10013", "EA216", "EA264", "EA364", "EA440", "EA741", "EA771", "EA876",
    "EA888", "EA987", "EA1054", "EA1081", "EA1161", "EA1167", "EA1233",
    "EA1302", "EA1506", "EA1511", "EA1513", "EA1581", "EA1630", "EA1885",
    "EA2067", "EA2249", "EA2254", "EA2256", "EA2304", "EA2372", "EA2389",
    "EA2403", "EA5682D", "EA5694", "EA5739", "EA5766", "EA5868", "EA5874",
    "EA5914", "EA5921", "EA5957", "EA6130", "EA6237", "EA6918", "EA7027",
    "EA7158", "EA8574", "EA9242", "EA9251", "EA9287", "EA9409", "EA9417",
    "EA9491", "EA9668", "EA9752", "EA9779", "EA11128", "PB47", "PB232",
    "PE171", "PT4-3",
]

COLORS = {
    "navy": "#17233E",
    "blue": "#1F77B4",
    "sky": "#5DADE2",
    "orange": "#F28E2B",
    "red": "#D62728",
    "green": "#2CA02C",
    "gray": "#7F8C8D",
    "light": "#F6F8FB",
    "grid": "#D9E2EC",
}

COLOR_CONDICION = {
    "Activo": COLORS["blue"],
    "Inactivo": COLORS["red"],
    "Observación": COLORS["orange"],
}

COLOR_TURNO = {
    "24 Hrs": COLORS["blue"],
    "12 Hrs": COLORS["orange"],
    "Convertido o CM": COLORS["green"],
    "Sin clasificar": COLORS["gray"],
}

st.set_page_config(page_title=APP_TITLE, page_icon="🛢️", layout="wide")


def inject_css():
    st.markdown(
        """
        <style>
        .main .block-container {padding-top: 1.4rem; padding-bottom: 2rem; max-width: 1500px;}
        .hero {
            background: linear-gradient(135deg, #17233E 0%, #1F77B4 60%, #2CA02C 100%);
            border-radius: 22px;
            padding: 28px 32px;
            color: white;
            margin-bottom: 18px;
            box-shadow: 0 14px 36px rgba(23,35,62,.20);
        }
        .hero h1 {font-size: 42px; line-height: 1.1; margin: 0; font-weight: 850;}
        .hero p {font-size: 16px; opacity: .92; margin: 10px 0 0 0;}
        .metric-card {
            background: white;
            border: 1px solid #E6EEF7;
            border-radius: 18px;
            padding: 16px 18px;
            box-shadow: 0 8px 22px rgba(23,35,62,.07);
        }
        .metric-label {font-size: 13px; color: #667085; margin-bottom: 4px;}
        .metric-value {font-size: 34px; color: #17233E; font-weight: 850; line-height: 1.05;}
        .metric-sub {font-size: 12px; color: #667085; margin-top: 6px;}
        .section-title {font-size: 24px; color: #17233E; font-weight: 850; margin: 18px 0 8px 0;}
        div[data-testid="stDownloadButton"] button {border-radius: 12px; font-weight: 700;}
        div[data-testid="stFormSubmitButton"] button {border-radius: 12px; font-weight: 800;}
        </style>
        """,
        unsafe_allow_html=True,
    )


def norm_txt(x) -> str:
    if pd.isna(x):
        return ""
    s = str(x).strip().upper()
    s = "".join(c for c in unicodedata.normalize("NFKD", s) if not unicodedata.combining(c))
    return re.sub(r"\s+", " ", s)


def pozo_key(x) -> str:
    """
    Limpia el código de pozo para comparar listas fijas.
    Quita espacios, guiones y símbolos no alfanuméricos.
    """
    s = norm_txt(x)
    return re.sub(r"[^A-Z0-9]", "", s)


POZOS_ATA_EXCLUIR_KEYS = {pozo_key(p) for p in POZOS_ATA_EXCLUIR}


def agregar_totales_barras_verticales(fig: go.Figure, totales: pd.DataFrame, x_col: str, y_col: str, formato="{:,.0f}") -> go.Figure:
    """
    Agrega etiquetas de total encima de barras verticales, especialmente en barras apiladas.
    """
    if totales.empty:
        return fig

    ymax = float(totales[y_col].max()) if y_col in totales.columns else 0
    offset = ymax * 0.035 if ymax > 0 else 1

    fig.add_trace(go.Scatter(
        x=totales[x_col],
        y=totales[y_col] + offset,
        mode="text",
        text=[formato.format(v) for v in totales[y_col]],
        textposition="top center",
        textfont=dict(size=13, color=COLORS["navy"], family="Arial Black"),
        showlegend=False,
        hoverinfo="skip",
    ))

    if ymax > 0:
        fig.update_yaxes(range=[0, ymax * 1.18])

    return fig


def obtener_columna_ubicacion(df: pd.DataFrame) -> str:
    """
    Para la lámina gerencial usa yacimiento si existe en el Excel.
    Si no existe, usa batería. En algunos archivos solo aparece SWAB como batería.
    """
    if "yacimiento_clean" in df.columns and df["yacimiento_clean"].replace("", np.nan).dropna().nunique() > 1:
        return "yacimiento_clean"
    return "bateria_clean"


def etiqueta_ubicacion(df: pd.DataFrame) -> str:
    col = obtener_columna_ubicacion(df)
    return "Yacimiento" if col == "yacimiento_clean" else "Batería / ubicación"



def mes_from_text(text: str):
    t = norm_txt(text)
    for m, n in MESES.items():
        if re.search(rf"\b{m}\b", t):
            return n
    return None


def year_from_text(text: str):
    m = re.search(r"\b(20\d{2}|19\d{2})\b", str(text))
    return int(m.group(1)) if m else None


def periodo_desde_nombre_archivo(file_name: str) -> pd.Timestamp | None:
    """
    Devuelve el periodo usando solo el nombre del archivo.

    Regla del proyecto:
    El mes y el año válidos son los que aparecen en el nombre del Excel,
    no los títulos internos ni el nombre de la hoja.
    """
    nombre = Path(file_name).stem
    mes = mes_from_text(nombre)
    anio = year_from_text(nombre)

    if mes is None or anio is None:
        return None

    return pd.Timestamp(int(anio), int(mes), 1)


def title_text(raw: pd.DataFrame) -> str:
    vals = []
    for v in raw.iloc[:4, :12].to_numpy().ravel():
        if pd.notna(v) and str(v).strip():
            vals.append(str(v))
    return " | ".join(vals)


def infer_period(file_name: str, sheet_name: str, raw: pd.DataFrame) -> pd.Timestamp | None:
    # Se conserva la firma por compatibilidad, pero la regla ahora es estricta:
    # el periodo sale únicamente del nombre del archivo.
    return periodo_desde_nombre_archivo(file_name)


def find_header_row(raw: pd.DataFrame) -> int | None:
    for i in range(min(len(raw), 15)):
        row = [norm_txt(x) for x in raw.iloc[i].tolist()]
        if all(c in row for c in REQ):
            return i
    return None


def extract_main_table(path: Path, sheet: str, header_row: int) -> pd.DataFrame:
    df = pd.read_excel(path, sheet_name=sheet, header=header_row, engine="openpyxl")
    clean_cols = [norm_txt(c) for c in df.columns]
    positions = {req: clean_cols.index(req) if req in clean_cols else None for req in REQ}
    if any(v is None for v in positions.values()):
        return pd.DataFrame()

    out = pd.DataFrame({
        "pozo": df.iloc[:, positions["*PFORMACION"]],
        "estado": df.iloc[:, positions["*ESTADO"]],
        "tipo_pozo": df.iloc[:, positions["*TIPO_DE_POZO"]],
        "ult_est": df.iloc[:, positions["*ULT_EST"]],
        "bateria": df.iloc[:, positions["*BATERIA"]],
    })

    # Si el Excel trae yacimiento/zona/área, se usa para la lámina gerencial.
    yac_nombres = [
        "*YACIMIENTO", "YACIMIENTO", "*YAC", "YAC",
        "*ZONA", "ZONA", "*AREA", "AREA",
        "*UBICACION", "UBICACION"
    ]
    yac_pos = next((i for i, c in enumerate(clean_cols) if c in yac_nombres), None)
    out["yacimiento"] = df.iloc[:, yac_pos] if yac_pos is not None else ""

    prev_cols = [i for i, c in enumerate(clean_cols) if c.startswith("*ULT_EST ") or c.startswith("*ULT_EST_")]
    if prev_cols:
        out["ult_est_anterior"] = df.iloc[:, prev_cols[0]]
    else:
        out["ult_est_anterior"] = ""

    out = out.dropna(subset=["pozo"])
    out = out[~out["pozo"].astype(str).str.startswith("*")].copy()
    for c in out.columns:
        out[c] = out[c].astype(str).str.strip().replace({"nan": "", "None": ""})
    out = out[out["pozo"].astype(str).str.strip() != ""]
    return out


def score_candidate(path: Path, sheet: str, periodo: pd.Timestamp, nrows: int) -> float:
    """
    Puntúa hojas candidatas dentro de un mismo Excel.
    El periodo NO sale de la hoja, solo del nombre del archivo.
    La hoja se elige por tener más filas útiles y, secundariamente,
    por parecer una hoja principal.
    """
    score = float(nrows)

    sheet_norm = norm_txt(sheet)
    if mes_from_text(sheet) == periodo.month:
        score += 1000
    if "HOJA" not in sheet_norm:
        score += 200
    if any(x in sheet_norm for x in ["ESTADO", "SWAB", "POZOS", "OIG"]):
        score += 300

    return score


def condicion_operativa(ult_est: str, estado: str, tipo: str) -> str:
    u = norm_txt(ult_est)
    e = norm_txt(estado)
    t = norm_txt(tipo)
    texto = f"{u} {e} {t}"
    if "INACT" in texto or any(w in texto for w in ["DPA", "ABAND", "SECO", "ATA"]):
        return "Inactivo"
    if "ACT" in texto or any(w in texto for w in ["PETS", "PETF", "BM", "BCP", "BO", "GL", "PL", "CM"]):
        return "Activo"
    return "Observación"


def es_swab(row) -> bool:
    texto = " ".join([
        norm_txt(row.get("bateria", "")),
        norm_txt(row.get("estado", "")),
        norm_txt(row.get("tipo_pozo", "")),
        norm_txt(row.get("ult_est", "")),
    ])
    return any(x in texto for x in ["SWAB", "SUAB", "PETS", "PETF"])


def turno_swab(row) -> str:
    texto = " ".join([norm_txt(row.get("estado", "")), norm_txt(row.get("tipo_pozo", "")), norm_txt(row.get("ult_est", ""))])
    if "CM" in texto:
        return "Convertido o CM"
    if "PETS" in texto or "SUABTBG" in texto:
        return "24 Hrs"
    if "PETF" in texto or "SUABCSG" in texto:
        return "12 Hrs"
    return "Sin clasificar"


def tipo_swab(row) -> str:
    turno = turno_swab(row)
    if turno == "24 Hrs":
        return "Básica 1"
    if turno == "12 Hrs":
        return "Básica 2"
    if turno == "Convertido o CM":
        return "Convertidos o CM"
    return "Otros SWAB"


def pretty_period(p: pd.Timestamp) -> str:
    return f"{MESES_INV[p.month]} {p.year}"


@st.cache_data(show_spinner="Cargando Excel fijos desde carpeta data...")
def cargar_base() -> Tuple[pd.DataFrame, pd.DataFrame]:
    if not DATA_DIR.exists():
        return pd.DataFrame(), pd.DataFrame()

    files = sorted(DATA_DIR.glob("Estado de Pozos*.xlsx"))
    candidates: Dict[pd.Timestamp, Dict] = {}
    log_rows = []

    for path in files:
        periodo_archivo = periodo_desde_nombre_archivo(path.name)

        if periodo_archivo is None:
            log_rows.append({
                "archivo": path.name,
                "hoja": "",
                "periodo": "",
                "filas": 0,
                "score": 0,
                "estado_carga": "Omitido: el nombre del archivo no tiene mes y año. Renombrar como 'Estado de Pozos - Abril 2026 OIG.xlsx'.",
            })
            continue

        try:
            xl = pd.ExcelFile(path, engine="openpyxl")
        except Exception as e:
            log_rows.append({
                "archivo": path.name,
                "hoja": "",
                "periodo": periodo_archivo.strftime("%Y-%m"),
                "filas": 0,
                "score": 0,
                "estado_carga": f"Error al abrir: {e}",
            })
            continue

        encontro_candidato = False

        for sheet in xl.sheet_names:
            try:
                raw = pd.read_excel(path, sheet_name=sheet, header=None, nrows=15, engine="openpyxl")
                header_row = find_header_row(raw)

                if header_row is None:
                    continue

                # El periodo se toma estrictamente del nombre del archivo.
                periodo = periodo_archivo

                df = extract_main_table(path, sheet, header_row)
                if df.empty:
                    continue

                encontro_candidato = True
                sc = score_candidate(path, sheet, periodo, len(df))

                log_rows.append({
                    "archivo": path.name,
                    "hoja": sheet,
                    "periodo": periodo.strftime("%Y-%m"),
                    "filas": len(df),
                    "score": round(sc, 3),
                    "estado_carga": "Candidato",
                })

                if periodo not in candidates or sc > candidates[periodo]["score"]:
                    candidates[periodo] = {
                        "df": df,
                        "score": sc,
                        "archivo": path.name,
                        "hoja": sheet,
                    }

            except Exception as e:
                log_rows.append({
                    "archivo": path.name,
                    "hoja": sheet,
                    "periodo": periodo_archivo.strftime("%Y-%m"),
                    "filas": 0,
                    "score": 0,
                    "estado_carga": f"Error hoja: {e}",
                })

        if not encontro_candidato:
            log_rows.append({
                "archivo": path.name,
                "hoja": "",
                "periodo": periodo_archivo.strftime("%Y-%m"),
                "filas": 0,
                "score": 0,
                "estado_carga": "Sin hoja válida: no se encontró la fila de encabezados requeridos.",
            })

    all_rows = []
    seleccionados = []

    for periodo, item in sorted(candidates.items()):
        df = item["df"].copy()
        df["fecha"] = periodo
        df["periodo"] = periodo.strftime("%Y-%m")
        df["mes"] = MESES_INV[periodo.month]
        df["anio"] = periodo.year
        df["archivo_fuente"] = item["archivo"]
        df["hoja_fuente"] = item["hoja"]
        all_rows.append(df)

        seleccionados.append({
            "archivo": item["archivo"],
            "hoja": item["hoja"],
            "periodo": periodo.strftime("%Y-%m"),
            "filas": len(df),
            "score": round(float(item["score"]), 3),
            "estado_carga": "Seleccionado",
        })

    log_rows.extend(seleccionados)

    if not all_rows:
        return pd.DataFrame(), pd.DataFrame(log_rows)

    base = pd.concat(all_rows, ignore_index=True)
    base["pozo_clean"] = base["pozo"].map(norm_txt)
    base["pozo_key"] = base["pozo"].map(pozo_key)

    # Excluir candidatos ATA antes de cualquier KPI, tabla o gráfica.
    base = base[~base["pozo_key"].isin(POZOS_ATA_EXCLUIR_KEYS)].copy()

    base["estado_clean"] = base["estado"].map(norm_txt)
    base["tipo_clean"] = base["tipo_pozo"].map(norm_txt)
    base["ult_est_clean"] = base["ult_est"].map(lambda x: norm_txt(x) or "SIN DATO")
    base["bateria_clean"] = base["bateria"].map(lambda x: norm_txt(x).replace("CA ", "CA-") or "SIN DATO")
    if "yacimiento" not in base.columns:
        base["yacimiento"] = ""
    base["yacimiento_clean"] = base["yacimiento"].map(lambda x: norm_txt(x) or "")
    base["es_swab"] = base.apply(es_swab, axis=1)
    base = base[base["es_swab"]].copy()
    base["condicion"] = base.apply(lambda r: condicion_operativa(r["ult_est"], r["estado"], r["tipo_pozo"]), axis=1)
    base["turno_swab"] = base.apply(turno_swab, axis=1)
    base["tipo_swab"] = base.apply(tipo_swab, axis=1)
    base["pozo_visible"] = base["pozo"].astype(str).str.strip()
    base["mes_label"] = base["fecha"].map(pretty_period)

    return base, pd.DataFrame(log_rows)


def base_mes(df: pd.DataFrame, fecha: pd.Timestamp) -> pd.DataFrame:
    return df[df["fecha"] == fecha].copy()


def n_pozos(df: pd.DataFrame) -> int:
    return int(df["pozo_clean"].nunique())


def pct(x: int | float, total: int | float) -> float:
    return 0.0 if total == 0 else float(x) / float(total) * 100


def metric_card(label: str, value: str, sub: str = ""):
    st.markdown(
        f"""
        <div class="metric-card">
            <div class="metric-label">{label}</div>
            <div class="metric-value">{value}</div>
            <div class="metric-sub">{sub}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def build_kpis(df: pd.DataFrame) -> Dict[str, object]:
    fechas = sorted(df["fecha"].unique())
    ini, fin = pd.Timestamp(fechas[0]), pd.Timestamp(fechas[-1])
    d_ini = base_mes(df, ini)
    d_fin = base_mes(df, fin)
    total_ini = n_pozos(d_ini)
    total_fin = n_pozos(d_fin)
    activos_fin = n_pozos(d_fin[d_fin["condicion"] == "Activo"])
    inactivos_fin = n_pozos(d_fin[d_fin["condicion"] == "Inactivo"])
    obs_fin = n_pozos(d_fin[d_fin["condicion"] == "Observación"])
    return {
        "ini": ini,
        "fin": fin,
        "total_ini": total_ini,
        "total_fin": total_fin,
        "delta_total": total_fin - total_ini,
        "activos_fin": activos_fin,
        "inactivos_fin": inactivos_fin,
        "obs_fin": obs_fin,
        "pct_activos": pct(activos_fin, total_fin),
        "pct_inactivos": pct(inactivos_fin, total_fin),
        "meses": len(fechas),
    }


def show_kpis(kpis: Dict[str, object]):
    c1, c2, c3, c4, c5 = st.columns(5)
    with c1:
        metric_card("Total Pozos Swab", f"{kpis['total_fin']:,}", f"Δ {kpis['delta_total']:+,} vs {pretty_period(kpis['ini'])}")
    with c2:
        metric_card("Activos", f"{kpis['activos_fin']:,}", f"{kpis['pct_activos']:.1f}% del cierre")
    with c3:
        metric_card("Inactivos", f"{kpis['inactivos_fin']:,}", f"{kpis['pct_inactivos']:.1f}% del cierre")
    with c4:
        metric_card("Observación", f"{kpis['obs_fin']:,}", "Registros no clasificados")
    with c5:
        metric_card("Meses analizados", f"{kpis['meses']:,}", f"{pretty_period(kpis['ini'])} a {pretty_period(kpis['fin'])}")


def layout_fig(fig: go.Figure, height: int = 470, legend: str = "") -> go.Figure:
    fig.update_layout(
        template="plotly_white",
        height=height,
        font=dict(family="Arial", color=COLORS["navy"], size=13),
        title=dict(font=dict(size=22, color=COLORS["navy"], family="Arial Black"), x=0.02),
        plot_bgcolor="white",
        paper_bgcolor="white",
        legend_title_text=legend,
        margin=dict(l=40, r=25, t=80, b=45),
    )
    fig.update_xaxes(showgrid=True, gridcolor=COLORS["grid"], zeroline=False)
    fig.update_yaxes(showgrid=True, gridcolor=COLORS["grid"], zeroline=False)
    return fig


def fig_estado_swab(df: pd.DataFrame) -> go.Figure:
    t = df.groupby(["fecha", "mes_label", "condicion"], as_index=False)["pozo_clean"].nunique()
    t = t.rename(columns={"pozo_clean": "pozos"})
    totales = df.groupby("fecha", as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "total"})

    fig = px.bar(
        t,
        x="fecha",
        y="pozos",
        color="condicion",
        text="pozos",
        color_discrete_map=COLOR_CONDICION,
        title="Estado de Pozos Swab por mes",
    )
    fig.update_traces(textposition="inside", textfont_size=12)
    fig = agregar_totales_barras_verticales(fig, totales, "fecha", "total")
    fig.update_xaxes(title="Mes", tickformat="%b %Y")
    fig.update_yaxes(title="Cantidad de pozos")
    return layout_fig(fig, 500, "Condición")

def fig_pct_inactivos(df: pd.DataFrame) -> go.Figure:
    t = df.groupby(["fecha", "condicion"], as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "pozos"})
    pivot = t.pivot(index="fecha", columns="condicion", values="pozos").fillna(0)
    pivot["Total"] = pivot.sum(axis=1)
    pivot["% Inactivo"] = np.where(pivot["Total"] > 0, pivot.get("Inactivo", 0) / pivot["Total"] * 100, 0)
    plot = pivot.reset_index()
    fig = go.Figure()
    fig.add_trace(go.Scatter(
        x=plot["fecha"],
        y=plot["% Inactivo"],
        mode="lines+markers+text",
        line=dict(width=4, color=COLORS["red"]),
        marker=dict(size=10),
        text=[f"{v:.1f}%" for v in plot["% Inactivo"]],
        textposition="top center",
        name="% Inactivo",
    ))
    fig.update_layout(title="Tendencia de pozos Swab inactivos")
    fig.update_xaxes(title="Mes", tickformat="%b %Y")
    fig.update_yaxes(title="% Inactivo", ticksuffix="%")
    return layout_fig(fig, 430)


def fig_turno_swab(df: pd.DataFrame) -> go.Figure:
    t = df.groupby(["fecha", "turno_swab"], as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "pozos"})
    fig = px.line(
        t,
        x="fecha",
        y="pozos",
        color="turno_swab",
        markers=True,
        text="pozos",
        color_discrete_map=COLOR_TURNO,
        title="Pozos Swab por tipo operativo",
    )
    fig.update_traces(line=dict(width=4), marker=dict(size=10), textposition="top center")
    fig.update_xaxes(title="Mes", tickformat="%b %Y")
    fig.update_yaxes(title="Cantidad de pozos")
    return layout_fig(fig, 470, "Tipo operativo")


def fig_tipo_cierre(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)

    t = (
        d.groupby(["tipo_swab", "condicion"], as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "pozos"})
    )
    totales = (
        d.groupby("tipo_swab", as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "total"})
    )

    fig = px.bar(
        t,
        x="tipo_swab",
        y="pozos",
        color="condicion",
        text="pozos",
        color_discrete_map=COLOR_CONDICION,
        title=f"Por Tipo de Pozo Swab al cierre de {pretty_period(fin)}",
    )
    fig.update_traces(textposition="inside", textfont_size=12)
    fig = agregar_totales_barras_verticales(fig, totales, "tipo_swab", "total")
    fig.update_xaxes(title="Tipo de pozo Swab")
    fig.update_yaxes(title="N° Pozos")
    return layout_fig(fig, 460, "Condición")


def fig_ubicacion_cierre(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    col_ubi = obtener_columna_ubicacion(d)
    nombre_ubi = etiqueta_ubicacion(d)

    t = (
        d.groupby([col_ubi, "condicion"], as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "pozos", col_ubi: "ubicacion"})
    )
    totales = (
        d.groupby(col_ubi, as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "total", col_ubi: "ubicacion"})
    )

    top = totales.sort_values("total", ascending=False).head(18)["ubicacion"]
    t = t[t["ubicacion"].isin(top)].copy()
    totales = totales[totales["ubicacion"].isin(top)].copy()

    orden = totales.sort_values("total", ascending=False)["ubicacion"].tolist()

    fig = px.bar(
        t,
        x="ubicacion",
        y="pozos",
        color="condicion",
        text="pozos",
        category_orders={"ubicacion": orden},
        color_discrete_map=COLOR_CONDICION,
        title=f"Por {nombre_ubi} al cierre de {pretty_period(fin)}",
    )
    fig.update_traces(textposition="inside", textfont_size=11)
    fig = agregar_totales_barras_verticales(fig, totales, "ubicacion", "total")
    fig.update_xaxes(title=nombre_ubi, tickangle=-45)
    fig.update_yaxes(title="N° Pozos")
    return layout_fig(fig, 520, "Condición")


def fig_lamina_ubicacion_tipo(df: pd.DataFrame) -> go.Figure:
    """
    Lámina gerencial similar a la referencia enviada:
    Total de pozos, distribución por ubicación y distribución por tipo de pozo.
    """
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    total_swab = n_pozos(d)
    activos = n_pozos(d[d["condicion"] == "Activo"])
    inactivos = n_pozos(d[d["condicion"] == "Inactivo"])

    col_ubi = obtener_columna_ubicacion(d)
    nombre_ubi = etiqueta_ubicacion(d)

    ubi = (
        d.groupby([col_ubi, "condicion"], as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "pozos", col_ubi: "ubicacion"})
    )
    ubi_tot = (
        d.groupby(col_ubi, as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "total", col_ubi: "ubicacion"})
        .sort_values("total", ascending=False)
        .head(14)
    )
    ubi = ubi[ubi["ubicacion"].isin(ubi_tot["ubicacion"])].copy()
    ubi_order = ubi_tot["ubicacion"].tolist()

    tipo = (
        d.groupby(["tipo_swab", "condicion"], as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "pozos"})
    )
    tipo_tot = (
        d.groupby("tipo_swab", as_index=False)["pozo_clean"]
        .nunique()
        .rename(columns={"pozo_clean": "total"})
        .sort_values("total", ascending=False)
    )
    tipo_order = tipo_tot["tipo_swab"].tolist()

    fig = make_subplots(
        rows=2,
        cols=1,
        vertical_spacing=0.20,
        subplot_titles=(f"Por {nombre_ubi}", "Por Tipo de Pozo")
    )

    # Ubicación
    for condicion in ["Activo", "Inactivo", "Observación"]:
        parte = ubi[ubi["condicion"] == condicion].set_index("ubicacion").reindex(ubi_order).fillna(0)
        fig.add_trace(
            go.Bar(
                x=ubi_order,
                y=parte["pozos"],
                name=condicion,
                marker_color=COLOR_CONDICION.get(condicion, COLORS["gray"]),
                text=[f"{int(v)}" if v > 0 else "" for v in parte["pozos"]],
                textposition="inside",
                legendgroup=condicion,
                showlegend=True,
            ),
            row=1,
            col=1
        )

    ymax1 = float(ubi_tot["total"].max()) if not ubi_tot.empty else 0
    fig.add_trace(
        go.Scatter(
            x=ubi_order,
            y=ubi_tot.set_index("ubicacion").reindex(ubi_order)["total"] + max(ymax1 * 0.035, 1),
            mode="text",
            text=[str(int(v)) for v in ubi_tot.set_index("ubicacion").reindex(ubi_order)["total"]],
            textfont=dict(size=12, color=COLORS["navy"], family="Arial Black"),
            showlegend=False,
            hoverinfo="skip",
        ),
        row=1,
        col=1
    )

    # Tipo
    for condicion in ["Activo", "Inactivo", "Observación"]:
        parte = tipo[tipo["condicion"] == condicion].set_index("tipo_swab").reindex(tipo_order).fillna(0)
        fig.add_trace(
            go.Bar(
                x=tipo_order,
                y=parte["pozos"],
                name=condicion,
                marker_color=COLOR_CONDICION.get(condicion, COLORS["gray"]),
                text=[f"{int(v)}" if v > 0 else "" for v in parte["pozos"]],
                textposition="inside",
                legendgroup=condicion,
                showlegend=False,
            ),
            row=2,
            col=1
        )

    ymax2 = float(tipo_tot["total"].max()) if not tipo_tot.empty else 0
    fig.add_trace(
        go.Scatter(
            x=tipo_order,
            y=tipo_tot.set_index("tipo_swab").reindex(tipo_order)["total"] + max(ymax2 * 0.035, 1),
            mode="text",
            text=[str(int(v)) for v in tipo_tot.set_index("tipo_swab").reindex(tipo_order)["total"]],
            textfont=dict(size=12, color=COLORS["navy"], family="Arial Black"),
            showlegend=False,
            hoverinfo="skip",
        ),
        row=2,
        col=1
    )

    fig.update_layout(
        title=f"Total de Pozos Swab: {total_swab:,} | Activos: {activos:,} | Inactivos: {inactivos:,} | Cierre: {pretty_period(fin)}",
        barmode="stack",
        template="plotly_white",
        height=900,
        font=dict(family="Arial", color=COLORS["navy"], size=13),
        title_font=dict(size=24, color=COLORS["navy"], family="Arial Black"),
        plot_bgcolor="white",
        paper_bgcolor="white",
        legend_title_text="Condición",
        margin=dict(l=45, r=30, t=95, b=60),
    )

    fig.update_xaxes(tickangle=-45, row=1, col=1)
    fig.update_yaxes(title="N° Pozos", row=1, col=1, range=[0, ymax1 * 1.18 if ymax1 else 10])
    fig.update_xaxes(row=2, col=1)
    fig.update_yaxes(title="N° Pozos", row=2, col=1, range=[0, ymax2 * 1.18 if ymax2 else 10])

    return fig

def fig_estado_final(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    t = d.groupby("ult_est_clean", as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "pozos"})
    t = t.sort_values("pozos", ascending=True).tail(14)
    fig = px.bar(
        t,
        x="pozos",
        y="ult_est_clean",
        orientation="h",
        text="pozos",
        color="pozos",
        color_continuous_scale="Blues",
        title=f"Estados operativos Swab al cierre de {pretty_period(fin)}",
    )
    fig.update_traces(textposition="outside", cliponaxis=False)
    fig.update_layout(coloraxis_showscale=False)
    fig.update_xaxes(title="Cantidad de pozos")
    fig.update_yaxes(title="Último estado")
    return layout_fig(fig, 520)


def fig_bateria_total(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    col_ubi = obtener_columna_ubicacion(d)
    nombre_ubi = etiqueta_ubicacion(d)

    t = d.groupby([col_ubi, "condicion"], as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "pozos", col_ubi: "ubicacion"})
    top = d.groupby(col_ubi)["pozo_clean"].nunique().sort_values(ascending=False).head(18).index
    t = t[t["ubicacion"].isin(top)]
    fig = px.bar(
        t,
        x="pozos",
        y="ubicacion",
        color="condicion",
        orientation="h",
        text="pozos",
        color_discrete_map=COLOR_CONDICION,
        title=f"Ubicación y condición de Pozos Swab al cierre de {pretty_period(fin)}",
    )
    fig.update_yaxes(categoryorder="total ascending", title=nombre_ubi)
    fig.update_xaxes(title="Cantidad de pozos")
    return layout_fig(fig, 620, "Condición")

def fig_bateria_inactivos(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    col_ubi = obtener_columna_ubicacion(d)
    nombre_ubi = etiqueta_ubicacion(d)

    total = d.groupby(col_ubi, as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "total", col_ubi: "ubicacion"})
    ina = d[d["condicion"] == "Inactivo"].groupby(col_ubi, as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "inactivos", col_ubi: "ubicacion"})
    t = total.merge(ina, on="ubicacion", how="left").fillna(0)
    t = t[t["total"] >= 5].copy()
    t["pct_inactivo"] = np.where(t["total"] > 0, t["inactivos"] / t["total"] * 100, 0)
    t = t.sort_values(["pct_inactivo", "inactivos"], ascending=True).tail(15)
    fig = px.bar(
        t,
        x="pct_inactivo",
        y="ubicacion",
        orientation="h",
        text=[f"{v:.1f}%" for v in t["pct_inactivo"]],
        color="pct_inactivo",
        color_continuous_scale="OrRd",
        title=f"{nombre_ubi}s con mayor proporción de pozos Swab inactivos en {pretty_period(fin)}",
    )
    fig.update_layout(coloraxis_showscale=False)
    fig.update_xaxes(title="% inactivo", ticksuffix="%")
    fig.update_yaxes(title=nombre_ubi)
    return layout_fig(fig, 540)

def fig_matriz_bateria_estado(df: pd.DataFrame) -> go.Figure:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    col_ubi = obtener_columna_ubicacion(d)
    nombre_ubi = etiqueta_ubicacion(d)

    top_ubi = d[col_ubi].value_counts().head(15).index
    top_est = d["ult_est_clean"].value_counts().head(10).index
    d = d[d[col_ubi].isin(top_ubi) & d["ult_est_clean"].isin(top_est)]
    pivot = d.pivot_table(index=col_ubi, columns="ult_est_clean", values="pozo_clean", aggfunc="nunique", fill_value=0)
    fig = px.imshow(
        pivot,
        text_auto=True,
        aspect="auto",
        color_continuous_scale="YlGnBu",
        title=f"Matriz {nombre_ubi} vs Estado Swab en {pretty_period(fin)}",
    )
    fig.update_xaxes(title="Último estado")
    fig.update_yaxes(title=nombre_ubi)
    return layout_fig(fig, 650)

def fig_variacion_mensual(df: pd.DataFrame) -> go.Figure:
    t = df.groupby("fecha", as_index=False)["pozo_clean"].nunique().rename(columns={"pozo_clean": "total_swab"})
    t["variacion"] = t["total_swab"].diff().fillna(0).astype(int)
    fig = px.bar(
        t,
        x="fecha",
        y="variacion",
        text="variacion",
        color="variacion",
        color_continuous_scale="RdYlGn",
        title="Variación mensual de Pozos Swab",
    )
    fig.update_layout(coloraxis_showscale=False)
    fig.update_traces(textposition="outside")
    fig.update_xaxes(title="Mes", tickformat="%b %Y")
    fig.update_yaxes(title="Variación de pozos")
    return layout_fig(fig, 430)


def cambios_entre_extremos(df: pd.DataFrame) -> pd.DataFrame:
    fechas = sorted(df["fecha"].unique())
    if len(fechas) < 2:
        return pd.DataFrame()
    ini, fin = fechas[0], fechas[-1]
    a = base_mes(df, ini)[["pozo_clean", "pozo_visible", "ult_est_clean", "condicion", "turno_swab", "tipo_swab", "bateria_clean"]].copy()
    b = base_mes(df, fin)[["pozo_clean", "pozo_visible", "ult_est_clean", "condicion", "turno_swab", "tipo_swab", "bateria_clean"]].copy()
    a = a.rename(columns={
        "ult_est_clean": "estado_inicial", "condicion": "condicion_inicial", "turno_swab": "turno_inicial",
        "tipo_swab": "tipo_inicial", "bateria_clean": "bateria_inicial", "pozo_visible": "pozo"
    })
    b = b.rename(columns={
        "ult_est_clean": "estado_final", "condicion": "condicion_final", "turno_swab": "turno_final",
        "tipo_swab": "tipo_final", "bateria_clean": "bateria_final", "pozo_visible": "pozo_final"
    })
    m = a.merge(b, on="pozo_clean", how="outer")
    m["pozo"] = m["pozo"].fillna(m["pozo_final"]).fillna(m["pozo_clean"])
    m["estado_inicial"] = m["estado_inicial"].fillna("NO ESTABA EN SWAB")
    m["estado_final"] = m["estado_final"].fillna("NO ESTÁ EN SWAB")
    m["condicion_inicial"] = m["condicion_inicial"].fillna("NO ESTABA EN SWAB")
    m["condicion_final"] = m["condicion_final"].fillna("NO ESTÁ EN SWAB")
    m["cambio_condicion"] = m["condicion_inicial"] + " → " + m["condicion_final"]
    m["cambio_estado"] = m["estado_inicial"] + " → " + m["estado_final"]
    m = m[(m["estado_inicial"] != m["estado_final"]) | (m["condicion_inicial"] != m["condicion_final"])]
    keep = ["pozo", "estado_inicial", "estado_final", "condicion_inicial", "condicion_final", "cambio_condicion", "cambio_estado", "bateria_inicial", "bateria_final", "tipo_inicial", "tipo_final"]
    return m[keep].sort_values(["cambio_condicion", "pozo"])


def fig_transiciones(df: pd.DataFrame) -> go.Figure | None:
    cambios = cambios_entre_extremos(df)
    if cambios.empty:
        return None
    t = cambios.groupby("cambio_condicion", as_index=False)["pozo"].nunique().rename(columns={"pozo": "pozos"})
    t = t.sort_values("pozos", ascending=True)
    fig = px.bar(
        t,
        x="pozos",
        y="cambio_condicion",
        orientation="h",
        text="pozos",
        color="pozos",
        color_continuous_scale="Blues",
        title="Transición de condición entre el primer y último mes seleccionado",
    )
    fig.update_layout(coloraxis_showscale=False)
    fig.update_traces(textposition="outside", cliponaxis=False)
    fig.update_xaxes(title="Cantidad de pozos")
    fig.update_yaxes(title="Transición")
    return layout_fig(fig, 460)


def fig_sankey_estado(df: pd.DataFrame) -> go.Figure | None:
    cambios = cambios_entre_extremos(df)
    if cambios.empty:
        return None
    t = cambios.groupby(["estado_inicial", "estado_final"], as_index=False)["pozo"].nunique().rename(columns={"pozo": "pozos"})
    t = t.sort_values("pozos", ascending=False).head(18)
    labels = list(pd.unique(t[["estado_inicial", "estado_final"]].values.ravel("K")))
    idx = {lab: i for i, lab in enumerate(labels)}
    fig = go.Figure(data=[go.Sankey(
        node=dict(label=labels, pad=18, thickness=18, color=COLORS["sky"]),
        link=dict(source=t["estado_inicial"].map(idx), target=t["estado_final"].map(idx), value=t["pozos"], color="rgba(31,119,180,0.25)"),
    )])
    fig.update_layout(title_text="Flujo de cambios de estado operativo", font_size=12)
    return layout_fig(fig, 560)


def conclusion_table(df: pd.DataFrame) -> pd.DataFrame:
    fin = df["fecha"].max()
    d = base_mes(df, fin)
    res = d.groupby("tipo_swab", as_index=False).agg(
        pozos=("pozo_clean", "nunique"),
        activos=("pozo_clean", lambda s: d.loc[s.index][d.loc[s.index, "condicion"] == "Activo"]["pozo_clean"].nunique()),
        inactivos=("pozo_clean", lambda s: d.loc[s.index][d.loc[s.index, "condicion"] == "Inactivo"]["pozo_clean"].nunique()),
    )
    res["% Activo"] = np.where(res["pozos"] > 0, res["activos"] / res["pozos"] * 100, 0).round(1)
    res["% Inactivo"] = np.where(res["pozos"] > 0, res["inactivos"] / res["pozos"] * 100, 0).round(1)
    total = pd.DataFrame({
        "tipo_swab": ["Total"],
        "pozos": [res["pozos"].sum()],
        "activos": [res["activos"].sum()],
        "inactivos": [res["inactivos"].sum()],
        "% Activo": [round(res["activos"].sum() / res["pozos"].sum() * 100, 1) if res["pozos"].sum() else 0],
        "% Inactivo": [round(res["inactivos"].sum() / res["pozos"].sum() * 100, 1) if res["pozos"].sum() else 0],
    })
    return pd.concat([res.sort_values("pozos", ascending=False), total], ignore_index=True)


def csv_bytes(df: pd.DataFrame) -> bytes:
    return df.to_csv(index=False, encoding="utf-8-sig").encode("utf-8-sig")



def safe_key(text: str) -> str:
    return re.sub(r"[^A-Za-z0-9_]+", "_", str(text)).strip("_").lower()

def fig_to_png(fig: go.Figure, width=1500, height=850) -> bytes:
    return fig.to_image(format="png", width=width, height=height, scale=2)


def add_header(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(23, 35, 62)
    bar.line.color.rgb = RGBColor(23, 35, 62)
    tx = slide.shapes.add_textbox(Inches(0.35), Inches(0.14), Inches(10.8), Inches(0.45))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    logo = slide.shapes.add_textbox(Inches(11.55), Inches(0.14), Inches(1.4), Inches(0.35))
    q = logo.text_frame.paragraphs[0]
    q.text = "OIG"
    q.alignment = PP_ALIGN.RIGHT
    q.font.size = Pt(18)
    q.font.bold = True
    q.font.color.rgb = RGBColor(255, 255, 255)
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.38), Inches(0.77), Inches(12.2), Inches(0.35))
        r = sb.text_frame.paragraphs[0]
        r.text = subtitle
        r.font.size = Pt(11)
        r.font.color.rgb = RGBColor(89, 97, 110)


def exportar_ppt(figs: List[Tuple[str, go.Figure]], df: pd.DataFrame, cambios: pd.DataFrame, resumen: pd.DataFrame) -> bytes:
    if Presentation is None:
        raise RuntimeError("No está instalado python-pptx")

    kpis = build_kpis(df)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    fondo = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = RGBColor(23, 35, 62)
    fondo.line.color.rgb = RGBColor(23, 35, 62)
    titulo = slide.shapes.add_textbox(Inches(0.72), Inches(1.65), Inches(11.5), Inches(1.2))
    p = titulo.text_frame.paragraphs[0]
    p.text = "Lote X"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = titulo.text_frame.add_paragraph()
    p2.text = "Pozos Swab"
    p2.font.size = Pt(52)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(Inches(0.78), Inches(3.55), Inches(11.5), Inches(0.6))
    q = sub.text_frame.paragraphs[0]
    q.text = f"Análisis mensual: {pretty_period(kpis['ini'])} a {pretty_period(kpis['fin'])}"
    q.font.size = Pt(21)
    q.font.color.rgb = RGBColor(255, 255, 255)

    slide = prs.slides.add_slide(blank)
    add_header(slide, "Resumen ejecutivo de Pozos Swab", f"Cierre: {pretty_period(kpis['fin'])}")
    metrics = [
        ("Total Pozos Swab", f"{kpis['total_fin']:,}", f"Δ {kpis['delta_total']:+,}"),
        ("Activos", f"{kpis['activos_fin']:,}", f"{kpis['pct_activos']:.1f}%"),
        ("Inactivos", f"{kpis['inactivos_fin']:,}", f"{kpis['pct_inactivos']:.1f}%"),
        ("Meses", f"{kpis['meses']:,}", "rango seleccionado"),
    ]
    for i, (label, value, subtxt) in enumerate(metrics):
        x = 0.55 + i * 3.15
        box = slide.shapes.add_shape(1, Inches(x), Inches(1.35), Inches(2.75), Inches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(246, 248, 251)
        box.line.color.rgb = RGBColor(217, 226, 236)
        tx = slide.shapes.add_textbox(Inches(x + 0.18), Inches(1.55), Inches(2.4), Inches(1.1))
        a = tx.text_frame.paragraphs[0]
        a.text = label
        a.font.size = Pt(11)
        a.font.color.rgb = RGBColor(89, 97, 110)
        b = tx.text_frame.add_paragraph()
        b.text = value
        b.font.size = Pt(30)
        b.font.bold = True
        b.font.color.rgb = RGBColor(23, 35, 62)
        c = tx.text_frame.add_paragraph()
        c.text = subtxt
        c.font.size = Pt(11)
        c.font.color.rgb = RGBColor(89, 97, 110)

    rows, cols = len(resumen) + 1, len(resumen.columns)
    table = slide.shapes.add_table(rows, cols, Inches(0.65), Inches(3.35), Inches(12.0), Inches(3.3)).table
    for j, col in enumerate(resumen.columns):
        table.cell(0, j).text = str(col)
    for i, row in resumen.iterrows():
        for j, col in enumerate(resumen.columns):
            table.cell(i + 1, j).text = str(row[col])
    for row in table.rows:
        for cell in row.cells:
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(9)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        for idx, (title, fig) in enumerate(figs, start=1):
            img = tmp / f"fig_{idx}.png"
            fig.write_image(str(img), width=1500, height=850, scale=2)
            slide = prs.slides.add_slide(blank)
            add_header(slide, title)
            slide.shapes.add_picture(str(img), Inches(0.45), Inches(0.95), width=Inches(12.45), height=Inches(6.15))

        if not cambios.empty:
            slide = prs.slides.add_slide(blank)
            add_header(slide, "Pozos Swab con cambio de estado", "Primeros 18 cambios detectados en el rango")
            top = cambios.head(18)[["pozo", "estado_inicial", "estado_final", "cambio_condicion"]]
            table = slide.shapes.add_table(len(top) + 1, 4, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.9)).table
            headers = ["Pozo", "Estado inicial", "Estado final", "Cambio condición"]
            for j, h in enumerate(headers):
                table.cell(0, j).text = h
            for i, (_, r) in enumerate(top.iterrows(), start=1):
                table.cell(i, 0).text = str(r["pozo"])
                table.cell(i, 1).text = str(r["estado_inicial"])
                table.cell(i, 2).text = str(r["estado_final"])
                table.cell(i, 3).text = str(r["cambio_condicion"])
            for row in table.rows:
                for cell in row.cells:
                    for par in cell.text_frame.paragraphs:
                        par.font.size = Pt(8)

        out = tmp / "Pozos_Swab_Lote_X_dashboard.pptx"
        prs.save(out)
        return out.read_bytes()


def ui():
    inject_css()
    st.markdown(
        """
        <div class="hero">
            <h1>Dashboard Pozos Swab Lote X</h1>
            <p>Base fija de Excel en carpeta data. El mes y año se toman del nombre del archivo. El análisis trabaja solamente con pozos Swab y se ejecuta al presionar el botón.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    base, log = cargar_base()
    if base.empty:
        st.error("No se encontró información válida. Verifique que exista la carpeta data o Data junto a app.py con los Excel de Estado de Pozos.")
        st.stop()

    periodos = sorted(base["fecha"].drop_duplicates().tolist())
    etiquetas = {p: pretty_period(pd.Timestamp(p)) for p in periodos}

    with st.sidebar:
        st.header("Panel de control")
        with st.form("form_swab"):
            inicio = st.selectbox("Mes inicial", periodos, index=0, format_func=lambda p: etiquetas[p])
            fin = st.selectbox("Mes final", periodos, index=len(periodos) - 1, format_func=lambda p: etiquetas[p])
            tipos = sorted(base["tipo_swab"].dropna().unique().tolist())
            condiciones = ["Activo", "Inactivo", "Observación"]
            tipo_sel = st.multiselect("Tipo de pozo Swab", tipos, default=tipos)
            condicion_sel = st.multiselect("Condición", condiciones, default=condiciones)
            baterias = sorted(base["bateria_clean"].dropna().unique().tolist())
            bateria_sel = st.multiselect("Baterías", baterias, default=[])
            ejecutar = st.form_submit_button("Ejecutar análisis", type="primary")

    with st.expander("Meses reconocidos y fuentes usadas", expanded=False):
        fuentes = base.groupby(["fecha", "archivo_fuente", "hoja_fuente"], as_index=False)["pozo_clean"].nunique()
        fuentes["periodo"] = fuentes["fecha"].dt.strftime("%Y-%m")
        st.dataframe(fuentes[["periodo", "archivo_fuente", "hoja_fuente", "pozo_clean"]].rename(columns={"pozo_clean": "pozos_swab"}), use_container_width=True)
        st.caption(f"Se excluyen del análisis {len(POZOS_ATA_EXCLUIR_KEYS)} pozos candidatos ATA.")
        if not log.empty:
            st.caption("Detalle de carga de archivos")
            st.dataframe(log, use_container_width=True, height=220)

    if not ejecutar and "ejecutado_swab" not in st.session_state:
        st.info("Seleccione el rango de meses y presione Ejecutar análisis.")
        st.stop()

    st.session_state["ejecutado_swab"] = True

    if inicio > fin:
        st.error("El mes inicial no puede ser mayor que el mes final.")
        st.stop()

    df = base[(base["fecha"] >= inicio) & (base["fecha"] <= fin)].copy()
    df = df[df["tipo_swab"].isin(tipo_sel) & df["condicion"].isin(condicion_sel)]
    if bateria_sel:
        df = df[df["bateria_clean"].isin(bateria_sel)]

    if df.empty:
        st.warning("No hay registros Swab con los filtros seleccionados.")
        st.stop()

    kpis = build_kpis(df)
    st.markdown(f"<div class='section-title'>Resumen del rango: {pretty_period(kpis['ini'])} a {pretty_period(kpis['fin'])}</div>", unsafe_allow_html=True)
    show_kpis(kpis)

    figs = [
        ("Lámina ubicación y tipo", fig_lamina_ubicacion_tipo(df)),
        ("Estado de Pozos Swab", fig_estado_swab(df)),
        ("Tendencia de inactivos", fig_pct_inactivos(df)),
        ("Tipo operativo 24 Hrs y 12 Hrs", fig_turno_swab(df)),
        ("Tipo de Pozos Swab", fig_tipo_cierre(df)),
        ("Ubicación y condición", fig_bateria_total(df)),
        ("Baterías con mayor inactividad", fig_bateria_inactivos(df)),
        ("Estados operativos al cierre", fig_estado_final(df)),
        ("Matriz ubicación vs Estado", fig_matriz_bateria_estado(df)),
        ("Variación mensual de Pozos Swab", fig_variacion_mensual(df)),
    ]
    trans = fig_transiciones(df)
    sankey = fig_sankey_estado(df)
    if trans is not None:
        figs.append(("Transición de condición", trans))
    if sankey is not None:
        figs.append(("Flujo de cambios de estado", sankey))

    resumen = conclusion_table(df)
    cambios = cambios_entre_extremos(df)

    tabs = st.tabs(["Resumen PPT", "Lámina ubicación y tipo", "Tendencias", "Ubicación", "Cambios", "Base", "Exportar"])

    with tabs[0]:
        c1, c2 = st.columns([1.1, 1])
        with c1:
            st.plotly_chart(figs[1][1], use_container_width=True, key="resumen_estado_swab")
        with c2:
            st.plotly_chart(figs[4][1], use_container_width=True, key="resumen_tipo_cierre")
        st.dataframe(resumen, use_container_width=True, hide_index=True)

    with tabs[1]:
        st.plotly_chart(figs[0][1], use_container_width=True, key="lamina_ubicacion_tipo")

    with tabs[2]:
        for i, (title, fig) in enumerate(figs[1:5]):
            st.plotly_chart(fig, use_container_width=True, key=f"tendencias_{i}_{safe_key(title)}")

    with tabs[3]:
        for i, (title, fig) in enumerate(figs[5:9]):
            st.plotly_chart(fig, use_container_width=True, key=f"ubicacion_{i}_{safe_key(title)}")

    with tabs[4]:
        if trans is not None:
            st.plotly_chart(trans, use_container_width=True, key="cambios_transicion_condicion")
        if sankey is not None:
            st.plotly_chart(sankey, use_container_width=True, key="cambios_sankey_estado")
        st.dataframe(cambios, use_container_width=True, height=430)

    with tabs[5]:
        st.dataframe(df, use_container_width=True, height=520)

    with tabs[6]:
        st.download_button("Descargar base Swab filtrada CSV", data=csv_bytes(df), file_name="base_swab_filtrada.csv", mime="text/csv")
        st.download_button("Descargar cambios CSV", data=csv_bytes(cambios), file_name="cambios_pozos_swab.csv", mime="text/csv")
        st.download_button("Descargar resumen CSV", data=csv_bytes(resumen), file_name="resumen_pozos_swab.csv", mime="text/csv")

        excluidos_df = pd.DataFrame({"pozos_ata_excluidos": sorted(POZOS_ATA_EXCLUIR_KEYS)})
        st.download_button(
            "Descargar lista de pozos ATA excluidos CSV",
            data=csv_bytes(excluidos_df),
            file_name="pozos_ata_excluidos.csv",
            mime="text/csv"
        )

        for title, fig in figs:
            safe = re.sub(r"[^A-Za-z0-9_]+", "_", title).strip("_").lower()
            try:
                st.download_button(
                    f"Descargar PNG: {title}",
                    data=fig_to_png(fig),
                    file_name=f"{safe}.png",
                    mime="image/png",
                )
            except Exception:
                pass

        try:
            ppt_bytes = exportar_ppt(figs, df, cambios, resumen)
            st.download_button(
                "Descargar PowerPoint estilo Pozos Swab",
                data=ppt_bytes,
                file_name="Pozos_Swab_Lote_X_dashboard.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
            )
        except Exception as e:
            st.warning(f"No se pudo generar la PPT. Revise kaleido y python-pptx en requirements.txt. Detalle: {e}")

if __name__ == "__main__":
    ui()
